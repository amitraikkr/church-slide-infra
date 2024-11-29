from pptx import Presentation
from pptx.util import Inches
from lxml import etree

def test_function():
    # Test lxml
    root = etree.Element("root")
    child = etree.SubElement(root, "child")
    child.text = "Hello, lxml!"
    print(etree.tostring(root, pretty_print=True).decode())

    # Test pptx
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello, pptx!"
    prs.save("test_presentation.pptx")
    print("Presentation saved as test_presentation.pptx")

if __name__ == "__main__":
    test_function()
