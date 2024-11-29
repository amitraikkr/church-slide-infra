import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.transitions import TransitionType
import os

s3 = boto3.client('s3')

# Get bucket names from environment variables
BUCKET_NAME = os.environ.get('BUCKET_NAME', 'churchfaithslideai-pptsv2')  # Fallback value for local testing
TEMPLATE_BUCKET = os.environ.get('TEMPLATE_BUCKET', 'slideai-templates')  # Fallback value for local testing

TEMPLATES = {
    'default': 'templates/default.potx',
    'modern': 'templates/modern.potx',
    'Christmas Service': 'templates/Christmas Service.potx',
    'minimal': 'templates/minimal.potx',
    'faith': 'templates/faith.potx'
}

# Add color mapping
THEME_COLORS = {
    'faith-purple': (108, 99, 255),  # #6C63FF
    'ocean-blue': (96, 165, 250),    # #60A5FA
    'emerald-green': (52, 211, 153), # #34D399
    'sunset-red': (248, 113, 113),   # #F87171
    'golden': (251, 191, 36)         # #FBBF24
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat'
}

# SLIDE_LAYOUTS = {
#     'title': 0,
#     'content': 1,
#     'section': 2,
#     'two_content': 3,
#     'comparison': 4,
#     'picture': 5
# }

def get_template(template_name='default'):
    """Get the PowerPoint template file based on template name"""
    try:
        template_key = TEMPLATES.get(template_name, TEMPLATES['default'])
        
        # Get template from S3
        template_obj = s3.get_object(Bucket=TEMPLATE_BUCKET, Key=template_key)
        template_buffer = io.BytesIO(template_obj['Body'].read())
        return Presentation(template_buffer)
    except Exception as e:
        print(f"Error loading template {template_name}: {str(e)}")
        return Presentation()  # Fallback to blank presentation

def create_title_slide(prs, title, theme):
    # Use title slide layout (usually index 0)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Get RGB values and font for the theme
    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
    
    # Add title
    title_shape = title_slide.shapes.title
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle if needed
    if title_slide.placeholders[1]:  # Subtitle placeholder
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Created with FaithSlide.ai"
        subtitle.text_frame.paragraphs[0].font.name = font
        subtitle.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def apply_theme(slide, theme):
    if not theme:
        return

    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')

    # Apply to title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.name = font
        title_frame.paragraphs[0].font.size = Pt(40)  # Updated size
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add bottom border to title
        title_shape = slide.shapes.title
        line = title_shape.line
        line.color.rgb = RGBColor(*rgb)
        line.width = Pt(2)

    # Apply to content
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = font
                paragraph.font.size = Pt(20)  # Updated size
                paragraph.font.color.rgb = RGBColor(*rgb)
                
                # Add proper spacing
                paragraph.space_before = Pt(12)
                paragraph.space_after = Pt(12)
                
                # Handle bullet points
                if paragraph.text.strip().startswith(('•', '-', '*')):
                    paragraph.level = 0
                    paragraph.bullet.style = 'BULLET'
                elif paragraph.text.strip().startswith(('1.', '2.', '3.')):
                    paragraph.level = 0
                    paragraph.bullet.style = 'ARABIC'

def add_slide_number(slide, index, total_slides):
    """Add slide number to the bottom right corner"""
    left = Inches(9)
    top = Inches(6.75)
    width = Inches(1)
    height = Inches(0.25)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"{index}/{total_slides}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

# def add_transition(slide):
#     slide.transition.type = TransitionType.MORPH

def handler(event, context):
    try:
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')
        
        # Get template name from theme
        template_name = theme.get('template', 'default')
        
        # Create presentation using specified template
        prs = get_template(template_name)
        
        # Create title slide
        create_title_slide(prs, presentation_title, theme)

        # Add content slides
        total_slides = len(slides) + 1  # +1 for title slide
        
        for index, slide_data in enumerate(slides, 1):
            # Choose appropriate layout based on content
            if '\n' in slide_data['content']:
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Text and content
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header
            
            slide.shapes.title.text = slide_data['title']
            slide.placeholders[1].text = slide_data['content']
            
            # Apply theme and add slide number
            apply_theme(slide, theme)
            add_slide_number(slide, index, total_slides - 1)

        # Save to memory buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"

        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'downloadUrl': url})
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'error': str(e)})
        }
