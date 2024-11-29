import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.transitions import TransitionType
import os

s3 = boto3.client('s3')

# Get bucket names from environment variables
BUCKET_NAME = os.environ.get('BUCKET_NAME', 'churchfaithslideai-pptsv2')  # Fallback value for local testing
TEMPLATE_BUCKET = os.environ.get('TEMPLATE_BUCKET', 'slideai-templates')  # Fallback value for local testing

TEMPLATES = {
    'default': 'templates/default.pptx',
    'modern': 'templates/modern.pptx',
    'christmas-service': 'templates/christmas-service.pptx',
    'minimal': 'templates/minimal.pptx',
    'faith': 'templates/faith.pptx'
}

# Add color mapping
THEME_COLORS = {
    'faith-purple': (108, 99, 255),  # #6C63FF
    'ocean-blue': (96, 165, 250),    # #60A5FA
    'emerald-green': (52, 211, 153), # #34D399
    'sunset-red': (248, 113, 113),   # #F87171
    'golden': (251, 191, 36)         # #FBBF24
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat'
}

# SLIDE_LAYOUTS = {
#     'title': 0,
#     'content': 1,
#     'section': 2,
#     'two_content': 3,
#     'comparison': 4,
#     'picture': 5
# }

def get_template(template_name='default'):
    """Get the PowerPoint template file based on template name"""
    try:
        # Convert template name to URL-safe format (replace spaces with hyphens)
        template_name = template_name.replace(' ', '-').lower()
        
        # Get template key, fallback to default if not found
        template_key = TEMPLATES.get(template_name, TEMPLATES['default'])
        print(f"Attempting to load template: {template_key} from bucket: {TEMPLATE_BUCKET}")
        
        # Get template from S3
        template_obj = s3.get_object(Bucket=TEMPLATE_BUCKET, Key=template_key)
        template_buffer = io.BytesIO(template_obj['Body'].read())
        prs = Presentation(template_buffer)
        
        # Print available layouts for debugging
        print("Available slide layouts:")
        for index, layout in enumerate(prs.slide_layouts):
            print(f"Layout {index}: {layout.name}")
            
        return prs
    except Exception as e:
        print(f"Error loading template {template_name}: {str(e)}")
        print(f"Falling back to default template")
        return Presentation()

def create_title_slide(prs, title, theme):
    """Create title slide using template's title layout"""
    # Get the title slide layout (usually the first layout)
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    
    # Find title and subtitle placeholders by their layout type
    title_placeholder = None
    subtitle_placeholder = None
    
    for shape in title_slide.placeholders:
        print(f"Shape type: {shape.placeholder_format.type}")  # Debug print
        if shape.placeholder_format.type == 1:  # Title
            title_placeholder = shape
        elif shape.placeholder_format.type == 2:  # Subtitle
            subtitle_placeholder = shape
    
    # Apply title and theme
    if title_placeholder:
        title_placeholder.text = title
        # Apply theme formatting
        for paragraph in title_placeholder.text_frame.paragraphs:
            paragraph.font.name = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
            paragraph.font.size = Pt(54)
            rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
            paragraph.font.color.rgb = RGBColor(*rgb)
    
    # Add subtitle
    if subtitle_placeholder:
        subtitle_placeholder.text = "Created with FaithSlide.ai"
        # Apply theme formatting
        for paragraph in subtitle_placeholder.text_frame.paragraphs:
            paragraph.font.name = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
            paragraph.font.size = Pt(32)
            rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
            paragraph.font.color.rgb = RGBColor(*rgb)

def create_content_slide(prs, slide_data, theme):
    """Create content slide using template's content layout"""
    # Use content layout (usually layout index 1 or 2)
    content_layout = prs.slide_layouts[1]  # Adjust index based on your template
    slide = prs.slides.add_slide(content_layout)
    
    # Find title and content placeholders
    title_placeholder = None
    content_placeholder = None
    
    for shape in slide.placeholders:
        print(f"Shape type: {shape.placeholder_format.type}")  # Debug print
        if shape.placeholder_format.type == 1:  # Title
            title_placeholder = shape
        elif shape.placeholder_format.type in [2, 15]:  # Content (body)
            content_placeholder = shape
    
    # Add title
    if title_placeholder:
        title_placeholder.text = slide_data['title']
        # Apply theme formatting
        for paragraph in title_placeholder.text_frame.paragraphs:
            paragraph.font.name = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
            paragraph.font.size = Pt(40)
            rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
            paragraph.font.color.rgb = RGBColor(*rgb)
    
    # Add content
    if content_placeholder:
        content_placeholder.text = slide_data['content']
        # Apply theme formatting
        for paragraph in content_placeholder.text_frame.paragraphs:
            paragraph.font.name = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
            paragraph.font.size = Pt(20)
            rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
            paragraph.font.color.rgb = RGBColor(*rgb)
            
            # Handle bullet points
            if paragraph.text.strip().startswith(('•', '-', '*')):
                paragraph.level = 0
                paragraph.bullet.style = 'BULLET'
            elif paragraph.text.strip().startswith(('1.', '2.', '3.')):
                paragraph.level = 0
                paragraph.bullet.style = 'ARABIC'
    
    return slide

def add_slide_number(slide, index, total_slides):
    """Add slide number to the bottom right corner"""
    left = Inches(9)
    top = Inches(6.75)
    width = Inches(1)
    height = Inches(0.25)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"{index}/{total_slides}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

# def add_transition(slide):
#     slide.transition.type = TransitionType.MORPH

def handler(event, context):
    try:
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')
        
        # Get template and create presentation
        prs = get_template(theme.get('template', 'default'))
        
        # Create title slide
        create_title_slide(prs, presentation_title, theme)
        
        # Add content slides
        total_slides = len(slides) + 1
        for index, slide_data in enumerate(slides, 1):
            slide = create_content_slide(prs, slide_data, theme)
            add_slide_number(slide, index, total_slides - 1)
        
        # Save to memory buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"

        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'downloadUrl': url})
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'error': str(e)})
        }
