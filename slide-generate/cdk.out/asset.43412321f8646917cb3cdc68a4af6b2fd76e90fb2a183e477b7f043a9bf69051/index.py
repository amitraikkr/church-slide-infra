import json
import boto3
import io
import os
import sys

# Wrap imports that might fail with better error handling
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as e:
    print(f"Failed to import python-pptx dependencies: {str(e)}")
    sys.exit(1)

s3 = boto3.client('s3')
BUCKET_NAME = 'churchfaithslideai-ppts'

def handler(event, context):
    try:
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        
        # Validate required fields
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')

        # Create presentation
        try:
            prs = Presentation()
        except Exception as e:
            raise RuntimeError(f"Failed to create presentation object: {str(e)}")

        # Add slides
        for index, slide_data in enumerate(slides):
            # Validate slide data
            if not isinstance(slide_data, dict):
                raise ValueError(f"Invalid slide data at index {index}")
            if 'title' not in slide_data or 'content' not in slide_data:
                raise ValueError(f"Slide at index {index} missing required 'title' or 'content'")

            try:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = slide_data['title']
                content = slide.placeholders[1]
                content.text = slide_data['content']
            except Exception as e:
                raise RuntimeError(f"Failed to create slide {index + 1}: {str(e)}")

        # Save to memory buffer
        try:
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
        except Exception as e:
            raise RuntimeError(f"Failed to save presentation to memory: {str(e)}")

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"
        
        # Upload to S3
        try:
            s3.put_object(
                Bucket=BUCKET_NAME,
                Key=file_name,
                Body=pptx_buffer.getvalue(),
                ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
        except Exception as e:
            raise RuntimeError(f"Failed to upload to S3: {str(e)}")

        # Generate presigned URL
        try:
            url = s3.generate_presigned_url(
                'get_object',
                Params={'Bucket': BUCKET_NAME, 'Key': file_name},
                ExpiresIn=3600
            )
        except Exception as e:
            raise RuntimeError(f"Failed to generate presigned URL: {str(e)}")

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({
                'downloadUrl': url
            })
        }

    except ValueError as e:
        print(f"Validation Error: {str(e)}")
        return {
            'statusCode': 400,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({
                'error': f'Validation Error: {str(e)}'
            })
        }
    except RuntimeError as e:
        print(f"Runtime Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({
                'error': f'Runtime Error: {str(e)}'
            })
        }
    except Exception as e:
        print(f"Unexpected Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({
                'error': 'An unexpected error occurred'
            })
        }