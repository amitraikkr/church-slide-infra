import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
import os

s3 = boto3.client('s3')
BUCKET_NAME = 'churchfaithslideai-ppts'

def handler(event, context):
    try:
        # Parse request body
        body = json.loads(event['body'])
        slides = body['slides']
        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')

        # Create presentation
        prs = Presentation()
        
        # Add slides
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using title and content layout
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data['title']
            
            # Set content
            content = slide.placeholders[1]
            content.text = slide_data['content']

        # Save to memory buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"
        
        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600  # URL expires in 1 hour
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({
                'downloadUrl': url
            })
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({
                'error': 'Failed to generate presentation'
            })
        }