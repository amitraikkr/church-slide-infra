import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.transitions import TransitionType
import os

s3 = boto3.client('s3')

# Get bucket names from environment variables
BUCKET_NAME = os.environ.get('BUCKET_NAME', 'churchfaithslideai-pptsv2')  # Fallback value for local testing
TEMPLATE_BUCKET = os.environ.get('TEMPLATE_BUCKET', 'slideai-templates')  # Fallback value for local testing

TEMPLATES = {
    'default': 'templates/default.pptx',
    'modern': 'templates/modern.pptx',
    'christmas-service': 'templates/christmas-service.pptx',
    'minimal': 'templates/minimal.pptx',
    'faith': 'templates/faith.pptx'
}

# Add color mapping
THEME_COLORS = {
    'faith-purple': (108, 99, 255),  # #6C63FF
    'ocean-blue': (96, 165, 250),    # #60A5FA
    'emerald-green': (52, 211, 153), # #34D399
    'sunset-red': (248, 113, 113),   # #F87171
    'golden': (251, 191, 36)         # #FBBF24
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat'
}

# SLIDE_LAYOUTS = {
#     'title': 0,
#     'content': 1,
#     'section': 2,
#     'two_content': 3,
#     'comparison': 4,
#     'picture': 5
# }

def get_template(template_name='default'):
    """Get the PowerPoint template file based on template name"""
    try:
        # Convert template name to URL-safe format (replace spaces with hyphens)
        template_name = template_name.replace(' ', '-').lower()
        
        # Get template key, fallback to default if not found
        template_key = TEMPLATES.get(template_name, TEMPLATES['default'])
        print(f"Attempting to load template: {template_key} from bucket: {TEMPLATE_BUCKET}")
        
        # Get template from S3
        template_obj = s3.get_object(Bucket=TEMPLATE_BUCKET, Key=template_key)
        template_buffer = io.BytesIO(template_obj['Body'].read())
        prs = Presentation(template_buffer)
        
        # Debug: Print available layouts
        print("Available slide layouts in template:")
        for idx, layout in enumerate(prs.slide_layouts):
            print(f"Layout {idx}: {layout.name}")
        
        return prs
    except Exception as e:
        print(f"Error loading template {template_name}: {str(e)}")
        print(f"Falling back to default template")
        return Presentation()

def apply_text_formatting(paragraph, font_name, font_size, rgb_color):
    """Apply consistent text formatting"""
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = RGBColor(*rgb_color)

def create_title_slide(prs, title, theme):
    """Create title slide using template's first layout"""
    # Get theme settings
    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
    
    # Add title slide using first layout
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Set title
    if title_slide.shapes.title:
        title_frame = title_slide.shapes.title.text_frame
        title_frame.text = title
        for paragraph in title_frame.paragraphs:
            apply_text_formatting(paragraph, font, 54, rgb)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Set subtitle
    for shape in title_slide.placeholders:
        if shape.placeholder_format.type == 2:  # Subtitle placeholder
            shape.text = "Created with FaithSlide.ai"
            for paragraph in shape.text_frame.paragraphs:
                apply_text_formatting(paragraph, font, 32, rgb)
                paragraph.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, slide_data, theme):
    """Create content slide using template layout"""
    # Get theme settings
    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
    
    # Choose layout based on content
    layout_index = 1 if '\n' in slide_data['content'] else 2
    print(f"Using layout index {layout_index} for slide with title: {slide_data['title']}")
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    
    # Debug: Print placeholders in the slide
    print(f"Placeholders in slide:")
    for shape in slide.placeholders:
        print(f"- Placeholder type: {shape.placeholder_format.type}, name: {shape.name}")
    
    # Set title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.text = slide_data['title']
        for paragraph in title_frame.paragraphs:
            apply_text_formatting(paragraph, font, 40, rgb)
            paragraph.alignment = PP_ALIGN.LEFT
    
    # Set content
    for shape in slide.placeholders:
        if shape.placeholder_format.type in [2, 15]:  # Content placeholder
            text_frame = shape.text_frame
            text_frame.text = slide_data['content']
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                apply_text_formatting(paragraph, font, 20, rgb)
                
                # Add proper spacing
                paragraph.space_before = Pt(12)
                paragraph.space_after = Pt(12)
                
                # Handle bullet points
                if paragraph.text.strip().startswith(('•', '-', '*')):
                    paragraph.level = 0
                    paragraph.bullet.style = 'BULLET'
                elif paragraph.text.strip().startswith(('1.', '2.', '3.')):
                    paragraph.level = 0
                    paragraph.bullet.style = 'ARABIC'

def add_slide_number(slide, index, total_slides):
    """Add slide number to the bottom right corner"""
    left = Inches(9)
    top = Inches(6.75)
    width = Inches(1)
    height = Inches(0.25)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"{index}/{total_slides}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

# def add_transition(slide):
#     slide.transition.type = TransitionType.MORPH

def handler(event, context):
    try:
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')
        
        # Get template name from theme
        template_name = theme.get('template', 'default')
        
        # Create presentation using template
        prs = get_template(template_name)
        
        # Create title slide
        create_title_slide(prs, presentation_title, theme)

        # Add content slides
        total_slides = len(slides) + 1  # +1 for title slide
        
        for index, slide_data in enumerate(slides, 1):
            slide = create_content_slide(prs, slide_data, theme)
            add_slide_number(slide, index, total_slides - 1)

        # Save presentation
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"

        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'downloadUrl': url})
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'error': str(e)})
        }
