import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

from io import BytesIO

s3 = boto3.client('s3')

# Get bucket name from environment variable
BUCKET_NAME = os.environ.get('BUCKET_NAME', 'churchfaithslideai-pptsv2')  # Fallback value for local testing
TEMPLATE_BUCKET = os.environ.get('TEMPLATE_BUCKET', 'slideai-templates')  # Fallback value for local testing
# Add color mapping
THEME_COLORS = {
    'faith-purple': (108, 99, 255),  # #6C63FF
    'ocean-blue': (96, 165, 250),    # #60A5FA
    'emerald-green': (52, 211, 153), # #34D399
    'sunset-red': (248, 113, 113),   # #F87171
    'golden': (251, 191, 36)         # #FBBF24
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat'
}

# Update background image paths to use S3
BACKGROUND_IMAGES = {
    'title': 'templates/christmas-service/title-bg.jpg',  # Title slide background
    'content': 'templates/christmas-service/content-bg.jpg'  # Content slide background
}

def get_background_image(image_key):
    """Get background image from S3"""
    try:
        print(f"Getting background image: {image_key}")
        response = s3.get_object(Bucket=TEMPLATE_BUCKET, Key=image_key)
        image_data = response['Body'].read()
        return BytesIO(image_data)
    except Exception as e:
        print(f"Error getting background image from S3: {str(e)}")
        return None

def apply_background_image(slide, image_key, prs, is_title=False):
    """Apply background image to slide with proper formatting"""
    try:
        # Get image from S3
        image_stream = get_background_image(image_key)
        if not image_stream:
            raise Exception("Failed to get background image from S3")
        
        # Add image to slide and ensure it's at the back
        left = top = 0
        width = prs.slide_width
        height = prs.slide_height
        
        pic = slide.shapes.add_picture(
            image_stream,
            left, top,
            width, height
        )
        
        # Ensure image is sent to back
        pic.z_order = 0
        
        print(f"Successfully applied background image: {image_key}")
        
    except Exception as e:
        print(f"Error applying background image: {str(e)}")
        # Fallback to solid color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)
        print("Applied fallback background color")

def create_title_slide(prs, title, theme):
    """Create title slide with custom formatting"""
    # Get theme settings
    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
    
    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Apply background image first
    apply_background_image(slide, BACKGROUND_IMAGES['title'], prs, is_title=True)
    
    # Add semi-transparent overlay with reduced transparency
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.fill.transparency = 0.6  # Reduced transparency to 60% for better visibility
    shape.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].font.bold = True
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(9), Inches(2)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Created with FaithSlide.ai"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.name = font
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

def create_content_slide(prs, slide_data, theme):
    """Create content slide with custom formatting"""
    # Get theme settings
    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')
    
    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Apply content background image from S3
    apply_background_image(slide, BACKGROUND_IMAGES['content'], prs)
    
    # Add semi-transparent overlay
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White overlay
    shape.fill.transparency = 0.85  # 85% transparent
    shape.line.fill.background()
    
    # Add title area with theme color
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(*rgb)
    title_bg.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add content area
    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.6),
        Inches(8.5), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Add content with formatting
    lines = slide_data['content'].split('\n')
    first = True
    for line in lines:
        if not first:
            paragraph = content_frame.add_paragraph()
        else:
            paragraph = content_frame.paragraphs[0]
            first = False
            
        paragraph.text = line.strip()
        paragraph.font.name = font
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)  # Dark grey text
        paragraph.space_before = Pt(12)
        paragraph.space_after = Pt(12)
        
        # Handle bullet points
        if line.strip().startswith(('•', '-', '*')):
            paragraph.level = 0
            paragraph.bullet.style = 'BULLET'
        elif line.strip().startswith(('1.', '2.', '3.')):
            paragraph.level = 0
            paragraph.bullet.style = 'ARABIC'
    
    return slide

def add_slide_number(slide, index, total_slides):
    """Add slide number to the bottom right corner"""
    left = Inches(9)
    top = Inches(6.75)
    width = Inches(1)
    height = Inches(0.25)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"{index}/{total_slides}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

def handler(event, context):
    try:
        # Add debug logging for background images
        print("Background image paths:")
        print(f"Title background: {BACKGROUND_IMAGES['title']}")
        print(f"Content background: {BACKGROUND_IMAGES['content']}")
        
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')
        
        # Create new presentation
        prs = Presentation()
        
        # Create title slide
        create_title_slide(prs, presentation_title, theme)

        # Add content slides
        total_slides = len(slides) + 1  # +1 for title slide
        
        for index, slide_data in enumerate(slides, 1):
            try:
                slide = create_content_slide(prs, slide_data, theme)
                if slide:  # Check if slide was created successfully
                    add_slide_number(slide, index, total_slides - 1)
                else:
                    print(f"Warning: Slide {index} was not created properly")
            except Exception as e:
                print(f"Error creating slide {index}: {str(e)}")
                continue  # Continue with next slide if one fails

        # Save presentation
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"

        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'downloadUrl': url})
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'error': str(e)})
        }
