import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

s3 = boto3.client('s3')
BUCKET_NAME = 'churchfaithslideai-pptsv2'

# Add color mapping
THEME_COLORS = {
    'faith-purple': (108, 99, 255),  # #6C63FF
    'ocean-blue': (96, 165, 250),    # #60A5FA
    'emerald-green': (52, 211, 153), # #34D399
    'sunset-red': (248, 113, 113),   # #F87171
    'golden': (251, 191, 36)         # #FBBF24
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat'
}

def apply_theme(slide, theme):
    if not theme:
        return

    # Get theme settings
    color_name = theme.get('color', 'faith-purple')
    font_name = theme.get('font', 'inter')
    
    # Get RGB values for the selected color
    rgb = THEME_COLORS.get(color_name, (108, 99, 255))  # Default to faith-purple
    
    # Get font name
    font = THEME_FONTS.get(font_name, 'Inter')  # Default to Inter

    # Apply to title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.name = font
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)

    # Apply to content
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = font
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = RGBColor(*rgb)

def handler(event, context):
    try:
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')

        # Create presentation
        prs = Presentation()

        # Add slides
        for index, slide_data in enumerate(slides):
            if 'title' not in slide_data or 'content' not in slide_data:
                raise ValueError(f"Slide at index {index} missing required 'title' or 'content'")
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data['title']
            slide.placeholders[1].text = slide_data['content']
            
            # Apply theme to the slide
            apply_theme(slide, theme)

        # Save to memory buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"

        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'downloadUrl': url})
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'error': str(e)})
        }
