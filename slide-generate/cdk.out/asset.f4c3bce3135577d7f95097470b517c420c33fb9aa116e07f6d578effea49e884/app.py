import json
import boto3
import io
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

s3 = boto3.client('s3')

# Environment variables for S3 buckets
BUCKET_NAME = os.environ.get('BUCKET_NAME', 'churchfaithslideai-pptsv2')
TEMPLATE_BUCKET = os.environ.get('TEMPLATE_BUCKET', 'slideai-templates')

TEMPLATES = {
    'default': 'templates/default.pptx',
    'modern': 'templates/modern.pptx',
    'christmas-service': 'templates/Christmas Service.pptx',
}

# Color and font mappings
THEME_COLORS = {
    'faith-purple': (108, 99, 255),
    'ocean-blue': (96, 165, 250),
    'emerald-green': (52, 211, 153),
    'sunset-red': (248, 113, 113),
    'golden': (251, 191, 36),
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat',
}

def get_template(template_name):
    """Fetch a PowerPoint template from S3."""
    template_key = TEMPLATES.get(template_name, TEMPLATES['default'])
    try:
        obj = s3.get_object(Bucket=TEMPLATE_BUCKET, Key=template_key)
        return Presentation(io.BytesIO(obj['Body'].read()))
    except Exception as e:
        raise ValueError(f"Error loading template '{template_name}': {e}")

def populate_slide_with_content(slide, title, content, theme):
    """Populate placeholders on the slide with title and content."""
    # Get theme colors and fonts
    rgb = THEME_COLORS.get(theme.get('color', 'faith-purple'), (108, 99, 255))
    font_name = THEME_FONTS.get(theme.get('font', 'inter'), 'Inter')

    # Set title text and style
    if slide.shapes.title:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        title_frame = title_placeholder.text_frame
        title_frame.paragraphs[0].font.name = font_name
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)

    # Set content text and style
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            text_frame = shape.text_frame
            text_frame.clear()  # Clear existing content
            for paragraph_text in content.split("\n"):
                paragraph = text_frame.add_paragraph()
                paragraph.text = paragraph_text
                paragraph.font.name = font_name
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(*rgb)

def create_presentation_with_template(template_name, slides_data, theme):
    """Create a presentation using a template and populate it with slide data."""
    # Load the template
    prs = get_template(template_name)

    # Add slides dynamically based on input data
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Default to 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get('title', 'Untitled Slide')
        content = slide_data.get('content', '')
        populate_slide_with_content(slide, title, content, theme)

    # Save the presentation to a buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer

def handler(event, context):
    """Lambda function entry point."""
    try:
        # Parse and validate input
        body = json.loads(event['body'])
        slides = body.get('slides', [])
        if not slides:
            raise ValueError("'slides' field is required and cannot be empty.")
        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Generated Presentation')
        template_name = theme.get('template', 'default')

        # Create the presentation
        pptx_buffer = create_presentation_with_template(template_name, slides, theme)

        # Define the output file key
        output_key = f"{presentation_title.replace(' ', '_')}.pptx"

        # Upload the generated presentation to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=output_key,
            Body=pptx_buffer.getvalue(),
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        # Generate a presigned URL for the client to download
        presigned_url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': output_key},
            ExpiresIn=3600  # URL expires in 1 hour
        )

        return {
            'statusCode': 200,
            'body': json.dumps({'downloadUrl': presigned_url})
        }

    except Exception as e:
        return {
            'statusCode': 500,
            'body': json.dumps({'error': str(e)})
        }
