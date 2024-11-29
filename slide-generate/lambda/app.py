import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

from io import BytesIO

s3 = boto3.client('s3')

# Get bucket name from environment variable
BUCKET_NAME = os.environ.get('BUCKET_NAME', 'churchfaithslideai-pptsv2')  # Fallback value for local testing
TEMPLATE_BUCKET = os.environ.get('TEMPLATE_BUCKET', 'slideai-templates')  # Fallback value for local testing
# Add color mapping
THEME_COLORS = {
    'faith-purple': (108, 99, 255),  # #6C63FF
    'ocean-blue': (96, 165, 250),    # #60A5FA
    'emerald-green': (52, 211, 153), # #34D399
    'sunset-red': (248, 113, 113),   # #F87171
    'golden': (251, 191, 36)         # #FBBF24
}

THEME_FONTS = {
    'inter': 'Inter',
    'georgia': 'Georgia',
    'poppins': 'Poppins',
    'playfair': 'Playfair Display',
    'montserrat': 'Montserrat'
}

# Add these constants at the top
TEMPLATE_PATHS = {
    'office-blue': 'templates/office-blue.pptx',
    'modern-dark': 'templates/modern-dark.pptx',
    'elegant-purple': 'templates/elegant-purple.pptx',
    'nature-green': 'templates/nature-green.pptx',
    'warm-coral': 'templates/warm-coral.pptx'
}

def get_template_presentation(theme_name):
    """Get the base template for the presentation"""
    template_key = TEMPLATE_PATHS.get(theme_name, 'templates/default.pptx')
    try:
        template_obj = s3.get_object(
            Bucket=TEMPLATE_BUCKET,
            Key=template_key
        )
        template_stream = io.BytesIO(template_obj['Body'].read())
        return Presentation(template_stream)
    except:
        # Fallback to blank presentation if template not found
        return Presentation()

def create_title_slide(prs, title, theme):
    """Create simple title slide with theme support"""
    # Get theme colors from the request
    background = theme.get('background', '#ffffff')
    title_color = theme.get('titleColor', '#000000')
    accent_color = theme.get('accentColor', '#6C63FF')
    font = theme.get('font', 'Calibri')
    
    # Convert hex colors to RGB tuples
    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    title_rgb = hex_to_rgb(title_color)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Set slide background if provided
    if background != '#ffffff':
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(background))
    
    # Add title text with theme styling
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_rgb)
    title_frame.paragraphs[0].font.bold = True

def create_content_slide(prs, slide_data, theme):
    """Create content slide with theme support"""
    # Get theme colors from the request
    background = theme.get('background', '#ffffff')
    title_color = theme.get('titleColor', '#000000')
    content_color = theme.get('contentColor', '#000000')
    font = theme.get('font', 'Calibri')
    
    # Convert hex colors to RGB
    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    title_rgb = hex_to_rgb(title_color)
    content_rgb = hex_to_rgb(content_color)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Set slide background if provided
    if background != '#ffffff':
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(background))
    
    # Add title with theme styling
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_rgb)
    title_frame.paragraphs[0].font.bold = True
    
    # Add content with theme styling
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8),
        Inches(8), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Format content with theme styling
    lines = slide_data['content'].split('\n')
    first = True
    for line in lines:
        if not first:
            paragraph = content_frame.add_paragraph()
        else:
            paragraph = content_frame.paragraphs[0]
            first = False
            
        paragraph.text = line.strip()
        paragraph.font.name = font
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(*content_rgb)
        paragraph.space_before = Pt(12)
        paragraph.space_after = Pt(12)
        
        if line.strip().startswith(('•', '-', '*')):
            paragraph.level = 0
            paragraph.bullet.style = 'BULLET'
        elif line.strip().startswith(('1.', '2.', '3.')):
            paragraph.level = 0
            paragraph.bullet.style = 'ARABIC'

    return slide

def add_slide_number(slide, index, total_slides):
    """Add slide number to the bottom right corner"""
    txBox = slide.shapes.add_textbox(
        Inches(9), Inches(6.75),
        Inches(0.5), Inches(0.25)
    )
    tf = txBox.text_frame
    tf.text = f"{index}/{total_slides}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

def apply_background_style(slide, theme):
    """Apply background style to slide"""
    if theme.get('template'):
        # If using a template, background is already set
        return
    
    background = theme.get('background', '#ffffff')
    if background.startswith('linear-gradient'):
        # Handle gradient backgrounds
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(*hex_to_rgb(theme['colors'][0]))
        fill.gradient_stops[1].color.rgb = RGBColor(*hex_to_rgb(theme['colors'][1]))
    elif background != '#ffffff':
        # Handle solid backgrounds
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(background))

def handler(event, context):
    try:
        # Validate input
        if not event.get('body'):
            raise ValueError("Request body is missing")

        # Parse request body
        body = json.loads(event['body'])
        if 'slides' not in body:
            raise ValueError("'slides' field is required in the request body")
        
        slides = body['slides']
        if not slides or not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")

        theme = body.get('theme', {})
        presentation_title = body.get('presentationTitle', 'Presentation')
        
        # Get theme information
        template_name = theme.get('template')

        # Create presentation from template
        prs = get_template_presentation(template_name)

        # Clear existing slides if using template
        if template_name:
            xml_slides = [slide._element for slide in prs.slides._sldIdLst]
            for slide in xml_slides:
                prs.slides._sldIdLst.remove(slide)

        # Create title slide
        title_slide = create_title_slide(prs, presentation_title, theme)
        apply_background_style(title_slide, theme)

        # Create content slides
        for index, slide_data in enumerate(slides, 1):
            content_slide = create_content_slide(prs, slide_data, theme)
            apply_background_style(content_slide, theme)
            add_slide_number(content_slide, index, len(slides))

        # Save presentation
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate unique filename
        file_name = f"{presentation_title.replace(' ', '_')}_{context.aws_request_id}.pptx"

        # Upload to S3
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=file_name,
            Body=pptx_buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Generate presigned URL
        url = s3.generate_presigned_url(
            'get_object',
            Params={
                'Bucket': BUCKET_NAME, 
                'Key': file_name,
                'ResponseContentType': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            },
            ExpiresIn=3600
        )

        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'downloadUrl': url})
        }

    except Exception as e:
        print(f"Error: {str(e)}")
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Credentials': True
            },
            'body': json.dumps({'error': str(e)})
        }
